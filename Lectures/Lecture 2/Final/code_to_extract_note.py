import os
from pptx import Presentation

def extract_slide_notes(presentation_path, output_folder=None):
    """
    Extract notes from each slide in a PowerPoint presentation and save as text files.
    
    Args:
        presentation_path (str): Path to the PowerPoint file
        output_folder (str): Folder to save text files (default: same as presentation)
    """
    
    # If no output folder specified, use the same folder as the presentation
    if output_folder is None:
        output_folder = os.path.dirname(presentation_path)
    
    # Create output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)
    
    try:
        # Load the presentation
        prs = Presentation(presentation_path)
        
        # Process each slide
        for i, slide in enumerate(prs.slides, 1):
            # Get the notes slide
            notes_slide = slide.notes_slide
            
            # Extract text from notes
            notes_text = ""
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text
            
            # Generate filename without zero-padding: Slide1.txt, Slide2.txt, etc.
            filename = f"Slide{i}.txt"
            filepath = os.path.join(output_folder, filename)
            
            # Write notes to file
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(notes_text)
            
            print(f"Created: {filename} - {len(notes_text)} characters")
        
        print(f"\nSuccessfully extracted notes from {len(prs.slides)} slides!")
        
    except Exception as e:
        print(f"Error: {e}")

def main():
    # Get the current directory
    current_dir = os.getcwd()
    
    # Look for PowerPoint files in the current directory
    pptx_files = [f for f in os.listdir(current_dir) if f.lower().endswith('.pptx')]
    
    if not pptx_files:
        print("No PowerPoint (.pptx) files found in the current directory.")
        return
    
    print("Found PowerPoint files:")
    for i, file in enumerate(pptx_files, 1):
        print(f"{i}. {file}")
    
    # If multiple files, let user choose
    if len(pptx_files) == 1:
        selected_file = pptx_files[0]
    else:
        try:
            choice = int(input("\nEnter the number of the file to process: ")) - 1
            selected_file = pptx_files[choice]
        except (ValueError, IndexError):
            print("Invalid selection.")
            return
    
    # Process the selected file
    presentation_path = os.path.join(current_dir, selected_file)
    print(f"\nProcessing: {selected_file}")
    
    extract_slide_notes(presentation_path)

if __name__ == "__main__":
    main()