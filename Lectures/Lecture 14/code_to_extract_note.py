import os
from pptx import Presentation

def extract_slide_notes(presentation_path: str, output_folder: str) -> None:
    """
    Extract notes from each slide in a PowerPoint presentation and save as text files.

    Args:
        presentation_path: Path to the PowerPoint file
        output_folder: Folder to save text files
    """
    os.makedirs(output_folder, exist_ok=True)

    try:
        prs = Presentation(presentation_path)

        for i, slide in enumerate(prs.slides, 1):
            # Notes may be missing
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text or ""

            filename = f"Slide{i}.txt"
            filepath = os.path.join(output_folder, filename)

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(notes_text)

            print(f"[OK] {os.path.basename(presentation_path)} -> {filename} ({len(notes_text)} chars)")

        print(f"Done: {os.path.basename(presentation_path)} ({len(prs.slides)} slides)\n")

    except Exception as e:
        print(f"[ERROR] {os.path.basename(presentation_path)}: {e}")

def main():
    # Base: wherever you run this script from
    base_dir = os.getcwd()

    # 1) Go to 'Final' folder
    final_dir = os.path.join(base_dir, "Final")
    if not os.path.isdir(final_dir):
        print(f"[ERROR] 'Final' folder not found at: {final_dir}")
        return

    # 4) Output root is ..\Texts (i.e., sibling of 'Final')
    output_root = os.path.join(base_dir, "Texts")
    os.makedirs(output_root, exist_ok=True)

    # 2) Pick all .pptx files in 'Final'
    pptx_files = [f for f in os.listdir(final_dir) if f.lower().endswith(".pptx")]
    if not pptx_files:
        print("[INFO] No .pptx files found in 'Final'.")
        return

    print(f"Found {len(pptx_files)} PowerPoint file(s) in {final_dir}:\n" +
          "\n".join(f" - {name}" for name in pptx_files) + "\n")

    for name in pptx_files:
        pres_path = os.path.join(final_dir, name)

        # 3) Keep naming of slide text files; also keep per-presentation folder name
        pres_stub = os.path.splitext(name)[0]
        pres_output_dir = os.path.join(output_root, pres_stub)

        extract_slide_notes(pres_path, pres_output_dir)

    print(f"All done. Text files saved under: {output_root}")

if __name__ == "__main__":
    main()
