import os
from pptx import Presentation

def extract_slide_notes(presentation_path: str, output_folder: str) -> None:
    """
    Extract notes from each slide in a PowerPoint presentation and save as text files
    directly inside output_folder, with names like Slide1.txt, Slide2.txt, etc.
    """
    os.makedirs(output_folder, exist_ok=True)

    try:
        prs = Presentation(presentation_path)

        for i, slide in enumerate(prs.slides, 1):
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text or ""

            filename = f"Slide{i}.txt"
            filepath = os.path.join(output_folder, filename)

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(notes_text)

            print(f"[OK] {os.path.basename(presentation_path)} -> {filename} ({len(notes_text)} chars)")

        print(f"Done: {os.path.basename(presentation_path)} ({len(prs.slides)} slides)\n")

    except Exception as e:
        print(f"[ERROR] {os.path.basename(presentation_path)}: {e}")

def main():
    base_dir = os.getcwd()
    final_dir = os.path.join(base_dir, "Final")

    if not os.path.isdir(final_dir):
        print(f"[ERROR] 'Final' folder not found at: {final_dir}")
        return

    output_root = os.path.join(base_dir, "Texts")
    os.makedirs(output_root, exist_ok=True)

    pptx_files = [f for f in os.listdir(final_dir) if f.lower().endswith(".pptx")]
    if not pptx_files:
        print("[INFO] No .pptx files found in 'Final'.")
        return

    print(f"Found {len(pptx_files)} PowerPoint file(s) in {final_dir}:\n" +
          "\n".join(f" - {name}" for name in pptx_files) + "\n")

    for name in pptx_files:
        pres_path = os.path.join(final_dir, name)
        extract_slide_notes(pres_path, output_root)

    print(f"All done. Text files saved under: {output_root}")

if __name__ == "__main__":
    main()
